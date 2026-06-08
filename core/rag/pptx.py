"""Parse .pptx into per-slide text. python-pptx is imported lazily so the module
is importable without the dependency installed."""

from dataclasses import dataclass


@dataclass
class SlideText:
    index: int  # 1-based slide number
    title: str
    body: str
    notes: str = ""


def parse_pptx(data: bytes) -> list[SlideText]:
    """Extract title / body / notes text per slide. Raises if python-pptx is
    missing or the bytes aren't a valid presentation."""
    from io import BytesIO

    from pptx import Presentation  # lazy: optional dependency

    prs = Presentation(BytesIO(data))
    slides: list[SlideText] = []
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_id = title_shape.shape_id if title_shape is not None else None
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
            if not text:
                continue
            if title_id is not None and shape.shape_id == title_id and not title:
                title = text
            else:
                body_parts.append(text)
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        body = "\n".join(body_parts).strip()
        # Many decks don't use the title *placeholder* — the heading is just the
        # first text box. When the placeholder is empty, fall back to the first
        # line of the body as the title and drop it from the body (so it isn't
        # duplicated once it's re-prefixed onto the chunk).
        if not title and body:
            head, _, rest = body.partition("\n")
            title, body = head.strip(), rest.strip()
        slides.append(SlideText(index=i, title=title, body=body, notes=notes))
    return slides
