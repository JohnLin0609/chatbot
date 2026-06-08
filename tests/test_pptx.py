"""parse_pptx: build a tiny presentation in-memory and parse it back."""

from io import BytesIO

from core.rag.pptx import parse_pptx


def _build_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content

    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "First Slide"
    s1.placeholders[1].text = "Body of the first slide"

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Second Slide"
    s2.placeholders[1].text = "More content here"

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_parse_pptx_per_slide():
    slides = parse_pptx(_build_pptx())
    assert len(slides) == 2
    assert slides[0].index == 1
    assert slides[0].title == "First Slide"
    assert "first slide" in slides[0].body.lower()
    assert slides[1].title == "Second Slide"


def _build_pptx_no_title_placeholder() -> bytes:
    """A deck whose heading lives in a plain text box (no title placeholder),
    mirroring the real course slides."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout: no title placeholder
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
    tf = tb.text_frame
    tf.text = "錯誤的種類"  # heading line
    tf.add_paragraph().text = "語法錯誤與執行時例外"  # content line

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_parse_pptx_falls_back_to_first_body_line_as_title():
    slides = parse_pptx(_build_pptx_no_title_placeholder())
    assert len(slides) == 1
    # no placeholder -> first line becomes the title, stripped from the body
    assert slides[0].title == "錯誤的種類"
    assert slides[0].body == "語法錯誤與執行時例外"
    assert "錯誤的種類" not in slides[0].body
