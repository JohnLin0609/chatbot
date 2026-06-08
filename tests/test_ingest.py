"""IngestService tests with fake embedding + vector store (doc_type="token" to
avoid loading spaCy in unit tests)."""

import pytest

from core.rag.ingest import IngestService, SlideRangeError, _lecture_from_filename
from core.tokens.counter import TokenCounter
from tests.conftest import make_settings

C = TokenCounter()


def _build_pptx(n: int) -> bytes:
    from io import BytesIO

    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    for i in range(1, n + 1):
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = f"Slide {i}"
        s.placeholders[1].text = f"body {i}"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


class FakeEmbedding:
    dim = 1536

    async def embed(self, texts):
        return [[float(i)] * 3 for i in range(len(texts))]


class FakeVectorStore:
    def __init__(self):
        self.deleted = []
        self.upserted = []

    async def delete_doc(self, doc_id):
        self.deleted.append(doc_id)

    async def upsert(self, points):
        self.upserted = points


def _svc(store):
    return IngestService(make_settings(ingest_chunk_tokens=20, ingest_chunk_overlap=4),
                         FakeEmbedding(), store, C)


async def test_ingest_short_doc_one_chunk():
    store = FakeVectorStore()
    doc_id, n = await _svc(store).ingest_text("a short note", title="Note", doc_type="token")
    assert n == 1
    assert len(store.upserted) == 1
    assert store.upserted[0].payload()["source"] == "curated"
    assert store.upserted[0].payload()["enabled"] is True
    assert store.upserted[0].title == "Note"


async def test_ingest_long_doc_multiple_chunks():
    store = FakeVectorStore()
    text = " ".join(f"word{i}" for i in range(200))
    doc_id, n = await _svc(store).ingest_text(text, doc_type="token")
    assert n > 1
    assert len(store.upserted) == n
    # delete-before-upsert keeps the doc clean
    assert store.deleted == [doc_id]


async def test_explicit_doc_id_respected():
    store = FakeVectorStore()
    doc_id, _ = await _svc(store).ingest_text("hi", doc_id="my-id", doc_type="token")
    assert doc_id == "my-id"


async def test_empty_text_no_chunks():
    store = FakeVectorStore()
    _doc_id, n = await _svc(store).ingest_text("   ", doc_type="token")
    assert n == 0
    assert store.upserted == []


async def test_ingest_pptx_skips_leading_and_trailing():
    store = FakeVectorStore()
    data = _build_pptx(4)  # slides 1..4
    _doc_id, n = await _svc(store).ingest_pptx(
        data, title="Deck", skip_leading=1, skip_trailing=1
    )
    assert n == 2
    # cover (1) + closing (4) dropped; slide_number provenance is preserved
    nums = [p.payload()["metadata"]["slide_number"] for p in store.upserted]
    assert nums == [2, 3]


async def test_ingest_pptx_skip_all_raises():
    store = FakeVectorStore()
    with pytest.raises(SlideRangeError):
        await _svc(store).ingest_pptx(_build_pptx(2), skip_leading=2)


def test_lecture_from_filename():
    assert _lecture_from_filename("W14_例外處理.pptx") == 14
    assert _lecture_from_filename("W05_條件判斷.py") == 5
    assert _lecture_from_filename("w7_for.py") == 7
    assert _lecture_from_filename("notes.txt") is None
    assert _lecture_from_filename(None) is None


async def test_ingest_pptx_stamps_content_type_and_lecture():
    store = FakeVectorStore()
    await _svc(store).ingest_pptx(_build_pptx(2), title="Deck",
                                  source_file="W08_期中複習.pptx")
    p = store.upserted[0]
    assert p.content_type == "slide" and p.lecture == 8
    assert p.source_file == "W08_期中複習.pptx" and p.language is None


async def test_ingest_code_one_chunk_with_fields():
    store = FakeVectorStore()
    svc = IngestService(make_settings(ingest_chunk_tokens=512), FakeEmbedding(), store, C)
    code = '"""第 5 週 範例"""\ndef pass_or_fail(s):\n    return s >= 60\n'
    doc_id, n = await svc.ingest_code(
        code, title="W05_條件判斷.py", source_file="W05_條件判斷.py", topic="conditionals")
    assert n == 1
    p = store.upserted[0]
    assert p.content_type == "code" and p.lecture == 5
    assert p.language == "python" and p.topic == "conditionals"
    assert p.source_file == "W05_條件判斷.py"
    assert "pass_or_fail" in p.text
